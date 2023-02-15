from pptx import Presentation
from pptx.util import Inches
from wand.image import Image

# Create a new PowerPoint presentation
presentation = Presentation()



# Load the logo image
with Image(filename='nike_black.png') as logo:
    logo.resize(1500,500)

    # Loop over the image files
    for image_file in ['image1.jpg', 'image2.jpg', 'image3.jpg', 'image4.jpg', 'image5.jpg']:
        # Load the image
        with Image(filename=image_file) as img:
            # Watermark the image with the logo
            img.composite(logo, 0, 0)
            # Save the watermarked image
            img.save(filename='output_' + image_file)
 
 
    # Loop over the image files
for i, image_file in enumerate(['output_image1.jpg', 'output_image2.jpg', 'output_image3.jpg', 'output_image4.jpg', 'output_image5.jpg']):
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    
    # Add the watermarked image to the slide
    
    picture = slide.shapes.add_picture(image_file, left=Inches(1), top=Inches(2.5), height=Inches(3), width=Inches(2))

    # Add the header and caption
    slide.shapes.title.text= "Sample Title {0}".format(i + 1)
    slide.placeholders[1].text = "Sample Subtitle {0}".format(i + 1)
 
# Save the PowerPoint presentation to the same directory
    presentation.save("presentation.pptx")